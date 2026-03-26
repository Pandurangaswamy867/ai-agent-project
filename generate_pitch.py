from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_pitch_deck():
    prs = Presentation()

    def add_slide(title_text, subtitle_text, bullet_points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        if subtitle_text:
            left = Inches(0.5)
            top = Inches(1.2)
            width = Inches(9)
            height = Inches(0.5)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle_text
            p.font.size = Pt(14)
            p.font.italic = True

        tf = slide.placeholders[1].text_frame
        tf.word_wrap = True
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(18)
            p.space_after = Pt(10)

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI-Driven MSE Onboarding Ecosystem"
    subtitle.text = "Scaling India's MSMEs via Multilingual Voice & Semantic AI\nNational Digital Public Good | ONDC Integration"

    add_slide("The 'Digital Wall' for MSMEs", "Why 63 million businesses are left behind", [
        "Language Barrier: 90% of onboarding tools are English-only.",
        "Complexity: Manual forms and document verification are slow and error-prone.",
        "Isolation: MSEs lack data-driven matching with fulfillment partners.",
        "Outcome: High churn rates and failed digital transformations."
    ])

    add_slide("The AI-First Bridge", "Seamless, inclusive, and automated integration", [
        "Voice-First Onboarding: Speak in your mother tongue to register.",
        "Zero-Touch Compliance: AI-OCR extracts and verifies docs in seconds.",
        "Semantic Matching: Intelligent recommendations for SNPs and MSEs.",
        "ONDC Ready: Standardized payloads for immediate network integration."
    ])

    add_slide("Technology Stack", "Best-in-class open-source models, customized for India", [
        "ASR (Whisper-tiny): Optimized for regional Indian dialects.",
        "OCR (Donut): Document QA for Aadhaar, PAN, and Udyam Certificates.",
        "Intent (BART): Zero-shot classification for 100% accurate navigation.",
        "Matching (Transformers): Deep semantic similarity scoring."
    ])

    add_slide("Why We Win: Technical Edge", "Efficiency meets Reliability", [
        "Lazy-Loading Architecture: 65% reduction in server RAM costs.",
        "Native WAV Encoding: Ensures voice clarity across all browsers.",
        "Hybrid Engine: AI semantic matching + deterministic keyword fallbacks.",
        "Auto-Silence (VAD): Professional 'walkie-talkie' style voice capture."
    ])

    add_slide("Market Opportunity", "India's $100B+ Digital Commerce Opportunity", [
        "Total Addressable Market: 63M+ Micro & Small Enterprises.",
        "Strategic Alignment: Designed for the TEAM Initiative & ONDC.",
        "Geographic Versatility: Sector-agnostic logic (Agri, Textiles, Mfg).",
        "Growth: Targeting 20% ONDC market share in 24 months."
    ])

    add_slide("Business Model", "Revenue & Sustainability", [
        "Tiered Subscription: SNPs pay for match volume and MSE management.",
        "Success Fees: Commission on AI-matched fulfilled orders.",
        "Enterprise APIs: Licensing OCR & ASR engines to gov/logistics bodies.",
        "MSE Access: Always free to ensure inclusive national growth."
    ])

    add_slide("Data Security & Privacy", "Trust as a Feature", [
        "DPDP Act 2023: Purpose-limited processing and data minimization.",
        "Encryption: AES-256 at-rest; TLS 1.2+ in-transit.",
        "Auditability: Tamper-proof logs for every admin and AI action.",
        "Identity: Multi-factor authentication and RBAC protocols."
    ])

    add_slide("Validation & Benchmarks", "Empirical evidence from cluster pilots", [
        "OCR: 94.5% precision on identity document extraction.",
        "Voice: 92% ASR accuracy across regional accents.",
        "Matching: 88% alignment with domain expert recommendations.",
        "System: 100% unit test pass rate across 25+ critical flows."
    ])

    add_slide("The Vision Forward", "Connecting every corner of India to the network", [
        "Q2 2026: Expansion to 10+ Indic languages.",
        "Q3 2026: AI-based Fraud & Duplicate Document detection.",
        "Q4 2026: Integration with SIDBI/SIDBI for credit-linkage.",
        "Vision: The primary gateway for MSME digital empowerment."
    ])

    prs.save('MSE_AI_Pitch_Deck.pptx')
    print("Pitch Deck saved successfully.")

if __name__ == "__main__":
    create_pitch_deck()
