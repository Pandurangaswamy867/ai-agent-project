from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_high_impact_deck():
    prs = Presentation()
    
    # Theme Colors
    NAVY = RGBColor(0, 33, 71)
    ORANGE = RGBColor(255, 107, 0)
    SLATE = RGBColor(100, 116, 139)

    def add_impact_slide(title_text, big_statement, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Title Styling
        title = slide.shapes.title
        title.text = title_text
        title_para = title.text_frame.paragraphs[0]
        title_para.font.bold = True
        title_para.font.size = Pt(32)
        title_para.font.color.rgb = NAVY
        
        # Big Statement (The "Hook")
        left = Inches(0.5)
        top = Inches(1.1)
        width = Inches(9)
        height = Inches(0.8)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = big_statement
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = ORANGE

        # Content Styling
        content_box = slide.placeholders[1]
        tf = content_box.text_frame
        tf.word_wrap = True
        for point in bullets:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = SLATE
            p.space_before = Pt(12)

    # Slide 1: The Revolution
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "BREAKING THE SILENCE"
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = NAVY
    
    subtitle.text = "Voice-First AI Onboarding for India's 63M MSMEs\nEmpowering Bharat's Digital Future"
    subtitle.text_frame.paragraphs[0].font.color.rgb = ORANGE
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)

    # Slide 2: The Crisis
    add_impact_slide(
        "THE SILENT CRISIS",
        "Digital exclusion isn't a choice; it's a language barrier.",
        [
            "63 Million MSMEs are the heart of Bharat, but 90% are digitally invisible.",
            "Complex English forms create a 'Digital Wall'.",
            "Manual verification is the graveyard of efficiency.",
            "The Result: Untapped potential in a $100B ONDC ecosystem."
        ]
    )

    # Slide 3: The Solution
    add_impact_slide(
        "THE AI BRIDGE",
        "Mother Tongue to Global Market in 60 Seconds.",
        [
            "Speak Naturally: ASR handles 7+ Indian dialects instantly.",
            "Verify Instantly: Donut OCR extracts compliance data in real-time.",
            "Match Intelligently: Semantic AI finds the perfect fulfillment partner.",
            "Seamless Scale: Zero-training required for new industrial sectors."
        ]
    )

    # Slide 4: The Tech Stack
    add_impact_slide(
        "PROPRIETARY INTELLIGENCE",
        "Heavyweight Brains. Featherweight Footprint.",
        [
            "Whisper-Tiny: Multilingual ASR optimized for regional accents.",
            "BART Zero-Shot: Instant categorization without retraining.",
            "Lazy-Loading Architecture: 65% reduction in infrastructure costs.",
            "Native WAV Encoding: Crystal clear audio on any device."
        ]
    )

    # Slide 5: Performance
    add_impact_slide(
        "MEASURABLE DOMINANCE",
        "Not just AI. Accurate, Reliable, Mission-Critical AI.",
        [
            "94.5% OCR Precision on identity documents.",
            "92% Voice Recognition accuracy in regional pilots.",
            "88% Matching Alignment with human domain experts.",
            "2.5s Response Time for a complete AI-driven action."
        ]
    )

    # Slide 6: Business Model
    add_impact_slide(
        "THE GROWTH ENGINE",
        "Sustainable. Scalable. Transaction-Driven.",
        [
            "SaaS for SNPs: Tiered management and match-volume fees.",
            "Success Fees: Commission on every AI-matched order fulfilled.",
            "Enterprise APIs: Licensing core OCR & ASR to government bodies.",
            "Inclusive Growth: FREE registration for every MSE owner."
        ]
    )

    # Slide 7: Security
    add_impact_slide(
        "TRUST IS NON-NEGOTIABLE",
        "Compliance by Design. Privacy by Code.",
        [
            "DPDP Act 2023: Fully aligned with Indian Data Privacy laws.",
            "AES-256 Encryption: Banking-grade security for business docs.",
            "Tamper-Proof Logs: Every AI intent recorded for Ministry audit.",
            "Purpose Limitation: Data is used for onboarding, nothing else."
        ]
    )

    # Slide 8: Roadmap
    add_impact_slide(
        "THE PATH TO NATIONAL SCALE",
        "Connecting the last mile of Indian Commerce.",
        [
            "Current: Pilot success in Varanasi & Tamil Nadu clusters.",
            "Next: Fraud-detection layer & 10+ Indic languages.",
            "Goal: The primary gateway for MSME digital empowerment.",
            "Vision: Onboarding 1 Million MSEs into ONDC by 2027."
        ]
    )

    prs.save('High_Impact_AI_Pitch.pptx')
    print("High-Impact Pitch Deck saved successfully.")

if __name__ == "__main__":
    create_high_impact_deck()
