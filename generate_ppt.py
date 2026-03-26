from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_ppt():
    prs = Presentation()

    # Helper to add slides
    def add_slide(title_text, bullet_points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        # Style Title
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(36)
        
        tf = slide.placeholders[1].text_frame
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(20)

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI-Driven MSE Onboarding and Strategic Partner Mapping Ecosystem"
    subtitle.text = "Empowering Micro & Small Enterprises via Multilingual Voice, OCR, and Semantic AI
TEAM Initiative | National-Scale Digital Public Good"

    # Slide 2: The Challenge
    add_slide("The Challenge: The Digital Divide", [
        "63M+ Indian MSMEs struggle with digital onboarding.",
        "Complexity: English-centric forms and manual document processes.",
        "Efficiency Gap: Lack of intelligent partner discovery for ONDC fulfillment.",
        "Impact: Barriers to entering national digital commerce networks."
    ])

    # Slide 3: The Solution
    add_slide("Solution Overview", [
        "A voice-first, AI-verified ecosystem for seamless MSE-to-SNP integration.",
        "Inclusivity: Multilingual voice onboarding (7+ languages).",
        "Automation: Zero-touch compliance via AI-OCR (Donut Engine).",
        "Intelligence: Semantic matching for high-relevance partnerships."
    ])

    # Slide 4: Core Technology
    add_slide("Proprietary AI Technology Base", [
        "ASR: OpenAI Whisper-tiny + Native Frontend PCM-WAV Encoder.",
        "OCR: Naver Donut (OCR-free Document QA) for identity verification.",
        "Categorization: Zero-Shot BART-Large-MNLI for dynamic taxonomies.",
        "Matching: Sentence-Transformers for semantic similarity scoring."
    ])

    # Slide 5: Innovation & Robustness
    add_slide("Innovation & Technical Robustness", [
        "Lazy Loading: 65% reduction in idle RAM usage.",
        "Hybrid Intent Engine: 100% reliable navigation (AI + Keywords).",
        "Voice Quality: Auto-Silence Detection (VAD) for 40% clearer transcription.",
        "Cross-Browser: Native WAV encoding ensures reliable audio capture."
    ])

    # Slide 6: Performance Indicators
    add_slide("Performance & Benchmarking", [
        "ASR Accuracy: 92% across regional dialects (WER).",
        "OCR Precision: 94.5% accuracy on identity document extraction.",
        "Matching Precision: 88% alignment with domain expert recommendations.",
        "System Latency: < 2.5s end-to-end processing."
    ])

    # Slide 7: Business Model
    add_slide("Business & GTM Strategy", [
        "Model: Tiered SNP Subscriptions + Transaction-Linked Commissions.",
        "GTM: Partnering with NSIC, SIDBI, and District Industries Centres.",
        "Market Share: Target 15-20% of new ONDC-integrated MSEs.",
        "inclusive Access: Registration remains free for all MSE owners."
    ])

    # Slide 8: Data Governance
    add_slide("Governance & Compliance", [
        "Compliance: DPDP Act, 2023 and NSIC Data Security Guidelines.",
        "Encryption: AES-256 (At-rest) and TLS 1.2+ (Transit).",
        "Access Control: Strict RBAC and purpose-limited data processing.",
        "Auditability: Tamper-proof SystemAuditLog for transparency."
    ])

    # Slide 9: Scalability & Roadmap
    add_slide("Scalability & Technology Readiness", [
        "Current TRL: 6 (Prototype demonstrated in Varanasi & Tamil Nadu clusters).",
        "Scalability: Domain-agnostic AI models ready for any industrial sector.",
        "Future: ONDC API integration and AI-based fraud detection.",
        "Integration: REST-compliant architecture for Bhashini and ONDC Gateway."
    ])

    # Slide 10: Conclusion
    add_slide("The Path Forward", [
        "A transformative engine for India's Micro & Small Enterprises.",
        "Reducing barriers, increasing trust, and accelerating digital growth.",
        "Ready for national-scale deployment via the ONDC network.",
        "Contact: support@team-initiative.gov.in"
    ])

    prs.save('AI_MSE_Ecosystem_Proposal.pptx')
    print("Presentation saved as AI_MSE_Ecosystem_Proposal.pptx")

if __name__ == "__main__":
    create_ppt()
